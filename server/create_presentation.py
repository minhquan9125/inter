"""
Script để tạo PowerPoint presentation cho 4 APIs
Yêu cầu: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Tạo presentation mới
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_title_slide(slide1)
    
    # Slide 2: Tổng quan
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_overview_slide(slide2)
    
    # Slide 3: API 1 - Export
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_api1_slide(slide3)
    
    # Slide 4: API 2 - Department
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_api2_slide(slide4)
    
    # Slide 5: API 3 - Specialization
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_api3_slide(slide5)
    
    # Slide 6: API 4 - Advanced Filters
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_api4_slide(slide6)
    
    # Slide 7: Authentication Flow
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_auth_slide(slide7)
    
    # Slide 8: Demo Instructions
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_demo_slide(slide8)
    
    # Slide 9: Comparison
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_comparison_slide(slide9)
    
    # Slide 10: Conclusion
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_conclusion_slide(slide10)
    
    # Lưu file
    prs.save('4_API_Quan_Ly_Bac_Si.pptx')
    print("✅ Đã tạo file: 4_API_Quan_Ly_Bac_Si.pptx")

def add_gradient_background(slide, color1, color2):
    """Thêm background cho slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color1)

def add_title_slide(slide):
    """Slide 1: Title"""
    add_gradient_background(slide, (102, 126, 234), (118, 75, 162))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "🏥 4 API QUẢN LÝ BÁC SĨ"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Hospital Management System"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(255, 215, 0)
    
    # Info
    info_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(1.5))
    info_frame = info_box.text_frame
    info_frame.text = "Ngày trình bày: 22/01/2026\n\nCông nghệ: Node.js + Express + MongoDB + JWT"
    for para in info_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(20)
        para.font.color.rgb = RGBColor(255, 255, 255)

def add_overview_slide(slide):
    """Slide 2: Tổng quan"""
    add_gradient_background(slide, (240, 248, 255), (240, 248, 255))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "📋 TỔNG QUAN"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(102, 126, 234)
    
    # Features grid
    features = [
        ("✅ 4 API Endpoints", "Export, Department,\nSpecialization, Filters"),
        ("🔐 JWT Authentication", "Token qua\nQuery String"),
        ("🎯 Flexible Filtering", "Multiple filters\nkết hợp"),
        ("📦 Export JSON", "Download toàn bộ\ndata")
    ]
    
    x_positions = [1, 5.5, 1, 5.5]
    y_positions = [2, 2, 4.5, 4.5]
    
    for i, (title, desc) in enumerate(features):
        box = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(x_positions[i]), Inches(y_positions[i]),
            Inches(3.5), Inches(2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(102, 126, 234)
        box.line.color.rgb = RGBColor(118, 75, 162)
        box.line.width = Pt(2)
        
        text_frame = box.text_frame
        text_frame.text = f"{title}\n\n{desc}"
        text_frame.margin_top = Inches(0.2)
        
        for para in text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            para.font.color.rgb = RGBColor(255, 255, 255)
            if para == text_frame.paragraphs[0]:
                para.font.size = Pt(20)
                para.font.bold = True
            else:
                para.font.size = Pt(16)

def add_api1_slide(slide):
    """Slide 3: API 1"""
    add_gradient_background(slide, (255, 255, 255), (255, 255, 255))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "📌 API 1: EXPORT JSON FILE"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(118, 75, 162)
    
    # API Box
    api_box = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.2), Inches(9), Inches(5.8)
    )
    api_box.fill.solid()
    api_box.fill.fore_color.rgb = RGBColor(248, 249, 250)
    api_box.line.color.rgb = RGBColor(102, 126, 234)
    api_box.line.width = Pt(4)
    
    text_frame = api_box.text_frame
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_top = Inches(0.2)
    
    content = """GET /api/doctors/export

Mô tả: Tải xuống toàn bộ danh sách bác sĩ dưới dạng JSON file
Token: Không cần ✅

URL: http://localhost:5000/api/doctors/export

Đặc điểm:
  ✅ Công khai - không yêu cầu authentication
  ✅ Trả về file JSON để download
  ✅ Chứa toàn bộ thông tin bác sĩ
  ✅ Dùng cho backup hoặc data export"""
    
    text_frame.text = content
    for i, para in enumerate(text_frame.paragraphs):
        if i == 0:  # Endpoint
            para.font.size = Pt(22)
            para.font.bold = True
            para.font.color.rgb = RGBColor(118, 75, 162)
        elif i in [4]:  # URL
            para.font.size = Pt(14)
            para.font.name = 'Courier New'
            para.font.color.rgb = RGBColor(72, 187, 120)
        else:
            para.font.size = Pt(18)
            para.font.color.rgb = RGBColor(51, 51, 51)

def add_api2_slide(slide):
    """Slide 4: API 2"""
    add_gradient_background(slide, (255, 255, 255), (255, 255, 255))
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "📌 API 2: GET BY DEPARTMENT"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(118, 75, 162)
    
    api_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    api_box.fill.solid()
    api_box.fill.fore_color.rgb = RGBColor(248, 249, 250)
    api_box.line.color.rgb = RGBColor(102, 126, 234)
    api_box.line.width = Pt(4)
    
    text_frame = api_box.text_frame
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_top = Inches(0.2)
    
    content = """GET /api/doctors/by-department

Mô tả: Lấy danh sách bác sĩ theo khoa/phòng ban
Token: Bắt buộc 🔐

URL: .../api/doctors/by-department?token=YOUR_TOKEN&department=Cardiology

Query Parameters:
  • token (string, required) - JWT authentication token
  • department (string, required) - Tên khoa (Cardiology, General...)

Ví dụ Department:
  ✅ Cardiology      ✅ Orthopedics
  ✅ General         ✅ Urology"""
    
    text_frame.text = content
    for i, para in enumerate(text_frame.paragraphs):
        if i == 0:
            para.font.size = Pt(22)
            para.font.bold = True
            para.font.color.rgb = RGBColor(118, 75, 162)
        elif i == 4:
            para.font.size = Pt(13)
            para.font.name = 'Courier New'
            para.font.color.rgb = RGBColor(72, 187, 120)
        else:
            para.font.size = Pt(16)
            para.font.color.rgb = RGBColor(51, 51, 51)

def add_api3_slide(slide):
    """Slide 5: API 3"""
    add_gradient_background(slide, (255, 255, 255), (255, 255, 255))
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "📌 API 3: GET BY SPECIALIZATION"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(118, 75, 162)
    
    api_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    api_box.fill.solid()
    api_box.fill.fore_color.rgb = RGBColor(248, 249, 250)
    api_box.line.color.rgb = RGBColor(102, 126, 234)
    api_box.line.width = Pt(4)
    
    text_frame = api_box.text_frame
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_top = Inches(0.2)
    
    content = """GET /api/doctors/by-specialization

Mô tả: Lấy danh sách bác sĩ theo chuyên khoa
Token: Bắt buộc 🔐

URL: .../by-specialization?token=YOUR_TOKEN&specialization=Cardiologist

Query Parameters:
  • token (string, required) - JWT authentication token
  • specialization (string, required) - Chuyên khoa

Ví dụ Specialization:
  ✅ Cardiologist         ✅ Orthopedic Surgeon
  ✅ Urologist            ✅ Gynaecologist"""
    
    text_frame.text = content
    for i, para in enumerate(text_frame.paragraphs):
        if i == 0:
            para.font.size = Pt(22)
            para.font.bold = True
            para.font.color.rgb = RGBColor(118, 75, 162)
        elif i == 4:
            para.font.size = Pt(13)
            para.font.name = 'Courier New'
            para.font.color.rgb = RGBColor(72, 187, 120)
        else:
            para.font.size = Pt(16)
            para.font.color.rgb = RGBColor(51, 51, 51)

def add_api4_slide(slide):
    """Slide 6: API 4"""
    add_gradient_background(slide, (255, 255, 255), (255, 255, 255))
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "📌 API 4: ADVANCED FILTERS"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(118, 75, 162)
    
    api_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    api_box.fill.solid()
    api_box.fill.fore_color.rgb = RGBColor(248, 249, 250)
    api_box.line.color.rgb = RGBColor(102, 126, 234)
    api_box.line.width = Pt(4)
    
    text_frame = api_box.text_frame
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_top = Inches(0.2)
    
    content = """GET /api/doctors

Mô tả: Lấy bác sĩ với nhiều bộ lọc kết hợp
Token: Bắt buộc 🔐

URL: .../doctors?token=TOKEN&name=Dr&specialization=...&department=...

Query Parameters:
  • token (string, required) - JWT authentication token
  • name (string, optional) - Tìm theo tên (case-insensitive)
  • specialization (string, optional) - Lọc theo chuyên khoa
  • department (string, optional) - Lọc theo khoa

Tính năng mạnh mẽ:
  ✅ Kết hợp nhiều filters cùng lúc
  ✅ Tìm kiếm không phân biệt hoa thường
  ✅ Trả về mảng rỗng nếu không tìm thấy"""
    
    text_frame.text = content
    for i, para in enumerate(text_frame.paragraphs):
        if i == 0:
            para.font.size = Pt(22)
            para.font.bold = True
            para.font.color.rgb = RGBColor(118, 75, 162)
        elif i == 4:
            para.font.size = Pt(12)
            para.font.name = 'Courier New'
            para.font.color.rgb = RGBColor(72, 187, 120)
        else:
            para.font.size = Pt(16)
            para.font.color.rgb = RGBColor(51, 51, 51)

def add_auth_slide(slide):
    """Slide 7: Authentication"""
    add_gradient_background(slide, (255, 255, 255), (255, 255, 255))
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "🔐 AUTHENTICATION FLOW"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(118, 75, 162)
    
    # Step 1
    step1_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.5), Inches(9), Inches(1.8))
    step1_box.fill.solid()
    step1_box.fill.fore_color.rgb = RGBColor(102, 126, 234)
    text_frame = step1_box.text_frame
    text_frame.text = """Bước 1: Login để lấy token
POST http://localhost:5000/api/auth/login
Body: {"email": "user@example.com", "password": "password"}"""
    for para in text_frame.paragraphs:
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.font.size = Pt(16)
        para.font.bold = True
    
    # Step 2
    step2_box = slide.shapes.add_shape(1, Inches(0.5), Inches(3.5), Inches(9), Inches(1.5))
    step2_box.fill.solid()
    step2_box.fill.fore_color.rgb = RGBColor(72, 187, 120)
    text_frame = step2_box.text_frame
    text_frame.text = """Bước 2: Sử dụng token trong URL
GET .../api/doctors/by-department?token=eyJhbGc..."""
    for para in text_frame.paragraphs:
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.font.size = Pt(16)
        para.font.bold = True
    
    # Middleware
    step3_box = slide.shapes.add_shape(1, Inches(0.5), Inches(5.2), Inches(9), Inches(1.8))
    step3_box.fill.solid()
    step3_box.fill.fore_color.rgb = RGBColor(248, 249, 250)
    step3_box.line.color.rgb = RGBColor(102, 126, 234)
    step3_box.line.width = Pt(2)
    text_frame = step3_box.text_frame
    text_frame.text = """Middleware isAuthenticated:
✅ Lấy token từ query string
✅ Verify JWT signature với secret key
✅ Extract user ID và role từ token
✅ Pass request nếu hợp lệ, reject nếu sai"""
    for para in text_frame.paragraphs:
        para.font.color.rgb = RGBColor(51, 51, 51)
        para.font.size = Pt(16)

def add_demo_slide(slide):
    """Slide 8: Demo"""
    add_gradient_background(slide, (255, 255, 255), (255, 255, 255))
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "🎯 HƯỚNG DẪN DEMO"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(118, 75, 162)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.7))
    text_frame = content_box.text_frame
    
    content = """1. Test API 1 (Export) - Công khai
   Mở browser: http://localhost:5000/api/doctors/export

2. Test API 2, 3, 4 - Cần Token

   Bước 1: Login bằng Postman
   POST /api/auth/login 
   Body: {"email":"user@example.com", "password":"password"}

   Bước 2: Copy token từ response

   Bước 3: Paste token vào URL
   GET /api/doctors/by-department?token=YOUR_TOKEN&department=Cardiology

💡 Tip: Import file Postman_Collection.json để test tự động!"""
    
    text_frame.text = content
    for i, para in enumerate(text_frame.paragraphs):
        if "1." in para.text or "2." in para.text or "Bước" in para.text:
            para.font.size = Pt(20)
            para.font.bold = True
            para.font.color.rgb = RGBColor(102, 126, 234)
        elif "Tip" in para.text:
            para.font.size = Pt(18)
            para.font.color.rgb = RGBColor(255, 140, 0)
            para.font.bold = True
        else:
            para.font.size = Pt(16)
            para.font.color.rgb = RGBColor(51, 51, 51)

def add_comparison_slide(slide):
    """Slide 9: Comparison"""
    add_gradient_background(slide, (255, 255, 255), (255, 255, 255))
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "📊 SO SÁNH: QUERY STRING vs HEADER"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(118, 75, 162)
    
    # Table content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    text_frame = content_box.text_frame
    
    content = """Tiêu chí          | Query String ✅      | Header
─────────────────────────────────────────────────────────
Test Browser      | Dễ dàng              | Không được
Share URL         | Copy paste đầy đủ    | Cần gửi token riêng
Debug             | Thấy ngay trong URL  | Phải mở DevTools
Bảo mật           | Token hiện trong URL | Ẩn trong header
Chuẩn REST        | Không theo chuẩn     | Theo chuẩn
─────────────────────────────────────────────────────────

✅ KẾT LUẬN: Query String phù hợp cho demo, testing, 
              internal APIs"""
    
    text_frame.text = content
    for para in text_frame.paragraphs:
        para.font.name = 'Courier New'
        if "KẾT LUẬN" in para.text:
            para.font.size = Pt(18)
            para.font.bold = True
            para.font.color.rgb = RGBColor(72, 187, 120)
        else:
            para.font.size = Pt(14)
            para.font.color.rgb = RGBColor(51, 51, 51)

def add_conclusion_slide(slide):
    """Slide 10: Conclusion"""
    add_gradient_background(slide, (102, 126, 234), (118, 75, 162))
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "✅ KẾT LUẬN"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    
    # Features
    features = [
        ("4 API Hoàn chỉnh", "Export, Department,\nSpecialization, Advanced Filters"),
        ("JWT Token", "Query String delivery,\n7 days expiration"),
        ("MongoDB Integration", "Mongoose ODM,\nIndexing, Regex search"),
        ("Postman Ready", "Collection file\nđầy đủ để test")
    ]
    
    x_positions = [1, 5.5, 1, 5.5]
    y_positions = [2, 2, 4.2, 4.2]
    
    for i, (title, desc) in enumerate(features):
        box = slide.shapes.add_shape(1, Inches(x_positions[i]), Inches(y_positions[i]), Inches(3.5), Inches(1.8))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        box.line.color.rgb = RGBColor(255, 215, 0)
        box.line.width = Pt(3)
        
        text_frame = box.text_frame
        text_frame.text = f"{title}\n\n{desc}"
        for para in text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            if para == text_frame.paragraphs[0]:
                para.font.size = Pt(18)
                para.font.bold = True
                para.font.color.rgb = RGBColor(102, 126, 234)
            else:
                para.font.size = Pt(14)
                para.font.color.rgb = RGBColor(51, 51, 51)
    
    # Thank you
    thank_box = slide.shapes.add_textbox(Inches(1), Inches(6.3), Inches(8), Inches(0.8))
    thank_frame = thank_box.text_frame
    thank_frame.text = "Cảm ơn thầy đã lắng nghe! 🙏"
    thank_para = thank_frame.paragraphs[0]
    thank_para.alignment = PP_ALIGN.CENTER
    thank_para.font.size = Pt(36)
    thank_para.font.bold = True
    thank_para.font.color.rgb = RGBColor(255, 215, 0)

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("❌ Lỗi: Chưa cài python-pptx")
        print("Chạy lệnh: pip install python-pptx")
    except Exception as e:
        print(f"❌ Lỗi: {e}")
